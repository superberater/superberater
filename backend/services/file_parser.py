"""File parser: Extract text from uploaded files.

Supported: PDF, DOCX, TXT, MD, XLSX, XLS, PPTX, CSV
"""

import io
import csv


async def extract_text(file_bytes: bytes, filename: str) -> str:
    """Extract text content from an uploaded file."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext in ("txt", "md"):
        return _extract_txt(file_bytes)
    elif ext == "csv":
        return _extract_csv(file_bytes)
    elif ext == "pdf":
        return _extract_pdf(file_bytes)
    elif ext in ("docx", "doc"):
        return _extract_docx(file_bytes)
    elif ext in ("xlsx", "xls"):
        return _extract_excel(file_bytes, ext)
    elif ext == "pptx":
        return _extract_pptx(file_bytes)
    else:
        supported = "txt, md, csv, pdf, docx, xlsx, xls, pptx"
        raise ValueError(f"Nicht unterstuetzter Dateityp: .{ext}. Unterstuetzt: {supported}")


def _extract_txt(data: bytes) -> str:
    for encoding in ("utf-8", "latin-1", "cp1252"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_csv(data: bytes) -> str:
    text = _extract_txt(data)
    reader = csv.reader(io.StringIO(text))
    rows = []
    for i, row in enumerate(reader):
        if i == 0:
            rows.append("| " + " | ".join(row) + " |")
            rows.append("| " + " | ".join(["---"] * len(row)) + " |")
        else:
            rows.append("| " + " | ".join(row) + " |")
        if i > 200:
            rows.append(f"... ({i}+ Zeilen, gekuerzt)")
            break
    return "\n".join(rows)


def _extract_pdf(data: bytes) -> str:
    try:
        import pdfplumber
    except ImportError:
        raise RuntimeError("pdfplumber nicht installiert. pip install pdfplumber")

    text_parts = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    return "\n\n".join(text_parts)


def _extract_docx(data: bytes) -> str:
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx nicht installiert. pip install python-docx")

    doc = Document(io.BytesIO(data))
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)
    return "\n\n".join(paragraphs)


def _extract_excel(data: bytes, ext: str) -> str:
    """Extract text from Excel files (xlsx/xls) using openpyxl."""
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("openpyxl nicht installiert. pip install openpyxl")

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    all_text = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        all_text.append(f"## Sheet: {sheet_name}")

        # Header
        header = rows[0]
        cols = [str(c) if c is not None else "" for c in header]
        all_text.append("| " + " | ".join(cols) + " |")
        all_text.append("| " + " | ".join(["---"] * len(cols)) + " |")

        # Data rows (max 200 per sheet)
        for i, row in enumerate(rows[1:], 1):
            cells = [str(c) if c is not None else "" for c in row]
            # Skip completely empty rows
            if not any(cells):
                continue
            all_text.append("| " + " | ".join(cells) + " |")
            if i > 200:
                all_text.append(f"... ({len(rows)}+ Zeilen, gekuerzt)")
                break

        all_text.append("")

    wb.close()
    return "\n".join(all_text)


def _extract_pptx(data: bytes) -> str:
    """Extract text from PowerPoint files using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx nicht installiert. pip install python-pptx")

    prs = Presentation(io.BytesIO(data))
    all_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells
                    )
                    if row_text.strip(" |"):
                        slide_texts.append(row_text)

        if slide_texts:
            all_text.append(f"--- Folie {slide_num} ---")
            all_text.append("\n".join(slide_texts))
            all_text.append("")

    return "\n".join(all_text)
